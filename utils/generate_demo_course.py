#!/usr/bin/env python3
"""Generate an AI-authored demo course and populate Canvas with it.

Pipeline: ask the configured LLM for a structured course outline
(app/course_gen.py) -> render one .pptx lecture deck per module with
python-pptx -> create the Canvas course, modules, and file attachments so
EasyLearn can immediately generate quizzes from the material.

Requires CANVAS_API_URL + CANVAS_API_TOKEN and a configured model provider
(GEMINI_API_KEY or OPENROUTER_API_KEY). See docs/demo.md.

Usage:
  uv run utils/generate_demo_course.py --topic "Introduction to Databases"
  uv run utils/generate_demo_course.py --topic "Astrophysics" --modules 4 --dry-run
  uv run utils/generate_demo_course.py --topic "Networking" --course-id 5 --enroll-teacher
"""

from __future__ import annotations

import argparse
import logging
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))
sys.path.insert(0, str(Path(__file__).resolve().parent))

from app import config  # noqa: E402
from app.canvas import get_canvas  # noqa: E402
from app.config import CACHE_DIR  # noqa: E402
from app.course_gen import generate_demo_course  # noqa: E402
from app.schemas import DemoModule  # noqa: E402

log = logging.getLogger(__name__)

DECK_DIR = CACHE_DIR / "demo_course"


def _slugify(text: str) -> str:
    keep = [c if c.isalnum() else "_" for c in text]
    return "".join(keep).strip("_")[:60] or "module"


def build_deck(module: DemoModule, dest_path: Path) -> Path:
    """Render a module's slides into a .pptx file at dest_path."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = module.name
    if module.summary and len(cover.placeholders) > 1:
        cover.placeholders[1].text = module.summary

    for slide in module.slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide.title
        body = s.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(slide.bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = bullet
            para.font.size = Pt(20)

    dest_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest_path))
    return dest_path


def upload_file(course, path: Path) -> int:
    ok, response = course.upload(str(path))
    if not ok:
        raise RuntimeError(f"Upload failed for {path.name}: {response}")
    file_id = response.get("id")
    if not file_id:
        raise RuntimeError(f"No file id in upload response for {path.name}: {response}")
    return int(file_id)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate an AI demo course into Canvas")
    parser.add_argument("--topic", required=True, help="Course topic to generate")
    parser.add_argument("--modules", type=int, default=3, help="Number of modules")
    parser.add_argument("--slides-per-module", type=int, default=6)
    parser.add_argument("--model-id", default=None, help="Model id from config/ai_models.json")
    parser.add_argument(
        "--course-id",
        type=int,
        default=None,
        help="Populate an existing course instead of creating one",
    )
    parser.add_argument(
        "--account-id",
        type=int,
        default=1,
        help="Canvas account id when creating a new course",
    )
    parser.add_argument(
        "--enroll-teacher",
        action="store_true",
        help="Also create/enroll the default test teacher (teacher1@example.com)",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Generate the outline and .pptx decks without touching Canvas",
    )
    parser.add_argument("-v", "--verbose", action="store_true")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s: %(message)s",
    )

    course_outline, model_entry = generate_demo_course(
        topic=args.topic,
        num_modules=args.modules,
        slides_per_module=args.slides_per_module,
        model_id=args.model_id,
    )
    print(f"Generated {len(course_outline.modules)} modules via {model_entry.label}")
    print(f"Course title: {course_outline.course_title}\n")

    decks: list[tuple[DemoModule, Path]] = []
    for position, module in enumerate(course_outline.modules, start=1):
        deck_path = DECK_DIR / f"{position:02d}_{_slugify(module.name)}.pptx"
        build_deck(module, deck_path)
        decks.append((module, deck_path))
        print(f"  [{position}] {module.name} ({len(module.slides)} slides) -> {deck_path.name}")

    if args.dry_run:
        print(f"\nDry run complete. Decks written to {DECK_DIR}")
        return 0

    canvas = get_canvas()

    if args.course_id:
        course = canvas.get_course(args.course_id)
        log.info("Using existing course id=%s: %s", course.id, course.name)
    else:
        account = canvas.get_account(args.account_id)
        course = account.create_course(
            course={
                "name": course_outline.course_title,
                "course_code": course_outline.course_code or "DEMO-101",
                "license": "private",
            }
        )
        course.update(course={"event": "offer"})
        log.info("Created course id=%s: %s", course.id, course.name)

    for position, (module, deck_path) in enumerate(decks, start=1):
        canvas_module = course.create_module({"name": module.name, "position": position})
        file_id = upload_file(course, deck_path)
        canvas_module.create_module_item(
            {"type": "File", "content_id": file_id, "title": deck_path.name}
        )
        canvas_module.edit(module={"published": True})
        log.info("  Module [%s]: %s (deck uploaded)", position, module.name)

    if args.enroll_teacher:
        try:
            from setup_canvas_test_users import DEFAULT_PASSWORD, enroll, ensure_user

            uid = ensure_user("teacher1@example.com", "Professor Ada", DEFAULT_PASSWORD)
            enroll(course.id, uid, "TeacherEnrollment")
            print(f"\nEnrolled teacher1@example.com (id={uid}) as Teacher")
        except Exception as exc:
            log.warning("Could not enroll teacher: %s", exc)

    print()
    print(f"Course URL: {config.CANVAS_API_URL}/courses/{course.id}")
    print(f"Modules created: {len(decks)}")
    print("Launch EasyLearn from this course in Canvas to generate quizzes.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
