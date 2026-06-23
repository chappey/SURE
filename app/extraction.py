"""Extract text from weekly course materials (PDF, PPTX)."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from .course_export import attachment_path

MAX_TEXT_CHARS = 100_000
MIN_TEXT_CHARS = 200

SUPPORTED_MATERIAL_SUFFIXES = (".pdf", ".pptx")


def is_supported_material(filename: str) -> bool:
    """Return True if the file type can be used for quiz generation."""
    lower = filename.lower()
    return lower.endswith(SUPPORTED_MATERIAL_SUFFIXES)


def extract_pdf_text(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts)


def extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = paragraph.text.strip()
                if line:
                    slide_parts.append(line)
        if slide_parts:
            parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_file_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    raise ValueError(f"Unsupported file type: {path.suffix} ({path.name})")


def extract_week_text(export_root: Path, module: dict[str, Any]) -> str:
    """Concatenate extracted text from all attachments in a week module."""
    sections: list[str] = []
    for item in module.get("items", []):
        if item.get("type") != "Attachment":
            continue
        path = attachment_path(export_root, item)
        text = extract_file_text(path).strip()
        header = f"## {path.name}"
        sections.append(f"{header}\n\n{text}" if text else f"{header}\n\n(no extractable text)")

    if not sections:
        raise ValueError(
            f"No attachments in module {module.get('name')!r}. "
            "Run course import or pick another week."
        )

    combined = "\n\n---\n\n".join(sections)
    if len(combined) > MAX_TEXT_CHARS:
        combined = combined[:MAX_TEXT_CHARS] + "\n\n[truncated]"
    return combined


def validate_week_text(text: str, week_name: str) -> None:
    if len(text.strip()) < MIN_TEXT_CHARS:
        raise ValueError(
            f"Extracted text for {week_name} is too short ({len(text)} chars). "
            "Check source files or extraction."
        )
